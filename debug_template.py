from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER

prs = Presentation("templates/base_template1.pptx")

for i, layout in enumerate(prs.slide_layouts):
    print(f"\nLayout {i}: {layout.name}")
    for shape in layout.placeholders:
        print(f"  - Placeholder {shape.placeholder_format.idx}: {shape.name} ({shape.placeholder_format.type})")
