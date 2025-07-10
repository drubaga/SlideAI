from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import PP_PLACEHOLDER, MSO_SHAPE
from pptx.dml.color import RGBColor
import os

from src.models.presentation import Presentation as PresentationModel

class PPTXBuilder:
    def __init__(
        self,
        presentation: PresentationModel,
        output_dir: str = "output",
        image_dir: str = "images",
        enable_images: bool = False,
        image_provider: str = None
    ):
        self.presentation_model = presentation
        self.output_dir = output_dir
        self.image_dir = image_dir
        self.enable_images = enable_images
        self.image_provider = image_provider
        self.template_name = presentation.template or "default"

        template_dir = os.path.abspath(os.path.join(os.path.dirname(__file__), "..", "..", "templates"))

        if self.template_name == "company":
            self.template_path = os.path.join(template_dir, "company_template.pptx")
        else:
            self.template_path = os.path.join(template_dir, "base_template1.pptx")

        self.prs = Presentation(self.template_path)

        # Set image fetcher if needed
        self.image_fetcher = None
        if self.enable_images and self.image_provider:
            if self.image_provider == "pexels":
                from src.image_generators.pexels_image_generator import PexelsImageGenerator
                self.image_fetcher = PexelsImageGenerator(self.image_dir)
            elif self.image_provider == "serpapi":
                from src.image_generators.serpapi_image_generator import SerpApiImageGenerator
                self.image_fetcher = SerpApiImageGenerator(self.image_dir)

    def build(self) -> str:
        if self.template_name == "company":
            self._remove_placeholder_middle_slide()
            self._insert_content_between_first_and_last()
        else:
            self._remove_all_initial_slides()
            self._add_title_slide()
            self._add_content_slides()
            self._add_thank_you_slide()
            self._remove_last_if_empty()


        return self._save_presentation()

    def _remove_last_if_empty(self):
        last_slide = self.prs.slides[-1]
        is_empty = True
        for shape in last_slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                is_empty = False
                break
        if is_empty:
            rId = self.prs.slides._sldIdLst[-1].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[-1]
    
    def _remove_all_initial_slides(self):
    # This avoids the need for corrupt _sldIdLst hacks
        while len(self.prs.slides):
            rId = self.prs.slides._sldIdLst[0].rId
            self.prs.part.drop_rel(rId)
            del self.prs.slides._sldIdLst[0]


    def _remove_placeholder_middle_slide(self):
        # Assuming slide 1 is the empty placeholder in company template
        if len(self.prs.slides) >= 3:
            xml_slides = self.prs.slides._sldIdLst
            del xml_slides[1]  # Remove 2nd slide (index 1)

    def _insert_content_between_first_and_last(self):
        layout = self.prs.slide_layouts[1]  # Assuming Title and Content
        insert_index = 1  # Between title and thank you

        for slide_data in self.presentation_model.slides:
            new_slide = self.prs.slides.add_slide(layout)
            self._fill_slide(new_slide, slide_data)

            # Insert newly created slide before thank you slide
            xml_slides = self.prs.slides._sldIdLst
            xml_slides.insert(insert_index, xml_slides[-1])
            insert_index += 1

    def _add_title_slide(self):
        layout = self.prs.slide_layouts[0]
        slide = self.prs.slides.add_slide(layout)
        slide.shapes.title.text = self.presentation_model.title
        if len(slide.placeholders) > 1:
            slide.placeholders[1].text = "An AI-generated presentation"

    def _add_content_slides(self):
        layout = self.prs.slide_layouts[1]
        for slide_data in self.presentation_model.slides:
            slide = self.prs.slides.add_slide(layout)
            self._fill_slide(slide, slide_data)

    def _add_thank_you_slide(self):
        layout = self.prs.slide_layouts[2]
        self.prs.slides.add_slide(layout)

    def _fill_slide(self, slide, slide_data):
        title_ph, content_ph = self._find_placeholders(slide)
        if not title_ph or not content_ph:
            raise ValueError("Missing title or content placeholder.")

        title_ph.text = slide_data.heading
        title_ph.text_frame.paragraphs[0].font.size = Pt(24)
        if self.template_name == "company":
            title_ph.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        text_frame = content_ph.text_frame
        text_frame.clear()

        for i, bullet in enumerate(slide_data.bullet_points or []):
            p = text_frame.add_paragraph() if i > 0 else text_frame.paragraphs[0]
            p.text = bullet
            p.level = 0
            p.font.size = Pt(18)
            if self.template_name == "company":
                p.font.color.rgb = RGBColor(255, 255, 255)

        if slide_data.key_message:
            p = text_frame.add_paragraph()
            p.text = slide_data.key_message
            p.level = 0
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = (
                RGBColor(255, 255, 255) if self.template_name == "company"
                else RGBColor(255, 105, 180)
            )

        if self.image_fetcher and slide_data.image_keywords:
            keywords = slide_data.image_keywords
            if isinstance(keywords, list):
                keywords = " ".join(keywords)
            image_path = self.image_fetcher.fetch_image(keywords)
            if image_path:
                self._insert_image_with_border(slide, image_path)

    def _find_placeholders(self, slide):
        title_ph, content_ph = None, None
        for shape in slide.placeholders:
            if shape.placeholder_format.type == PP_PLACEHOLDER.TITLE:
                title_ph = shape
            elif shape.placeholder_format.type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT):
                content_ph = shape
        return title_ph, content_ph

    def _insert_image_with_border(self, slide, image_path):
        margin = Inches(0.3)
        width, height = Inches(4.0), Inches(3.0)
        left = self.prs.slide_width - width - margin
        top = self.prs.slide_height - height - margin

        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left - Pt(2), top - Pt(2), width + Pt(4), height + Pt(4)
        )
        slide.shapes._spTree.remove(border._element)
        slide.shapes._spTree.insert(2, border._element)
        border.fill.background()
        border.line.width = Pt(2)
        border.line.color.rgb = RGBColor(0, 0, 0)

        slide.shapes.add_picture(image_path, left, top, width=width, height=height)

    def _save_presentation(self) -> str:
        os.makedirs(self.output_dir, exist_ok=True)
        output_path = os.path.join(self.output_dir, f"{self.presentation_model.title.replace(' ', '_')}.pptx")
        self.prs.save(output_path)
        return output_path
